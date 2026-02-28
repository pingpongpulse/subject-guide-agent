from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
import pdfplumber
import docx
from pptx import Presentation
from document_processor.ocr_loader import load_image, load_scanned_pdf
import os

def extract_text_from_pdf(file_path: str) -> list[dict]:
    """Extract digital text page by page from PDF."""
    pages = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    pages.append({
                        "text": text,
                        "page_number": i + 1
                    })
    except Exception as e:
        print(f"⚠️ Error reading digital PDF {file_path}: {e}")
    return pages

def extract_text_from_docx(file_path: str) -> list[dict]:
    """Extract text from DOCX."""
    doc = docx.Document(file_path)
    full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return [{"text": full_text, "page_number": 1}]

def extract_text_from_pptx(file_path: str) -> list[dict]:
    """Extract text slide by slide from PPTX."""
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
        if text.strip():
            slides.append({
                "text": text,
                "page_number": i + 1
            })
    return slides

def detect_doc_type(file_name: str) -> str:
    """Categorize document based on filename."""
    name_lower = file_name.lower()
    if any(k in name_lower for k in ["pyq", "exam", "paper"]):
        return "pyq"
    if any(k in name_lower for k in ["slide", "ppt", "lecture"]):
        return "lecture_slides"
    return "notes"

def chunk_document(file_path: str, subject: str = "general") -> list[Document]:
    """
    Main entry point. Loads file, handles OCR for images/scanned PDFs, 
    and splits into chunks for ChromaDB.
    """
    file_name = os.path.basename(file_path)
    extension = file_name.split(".")[-1].lower()
    doc_type = detect_doc_type(file_name)

    # --- STEP 1: Load Content ---
    pages = []
    
    if extension == "pdf":
        # First try standard extraction, if empty, use OCR
        pages = extract_text_from_pdf(file_path)
        if not pages:
            print(f"🔍 PDF {file_name} looks scanned. Switching to OCR...")
            pages = load_scanned_pdf(file_path)
            
    elif extension in ["png", "jpg", "jpeg"]:
        print(f"📸 Processing Image with OCR: {file_name}")
        pages = load_image(file_path)
        
    elif extension == "docx":
        pages = extract_text_from_docx(file_path)
        
    elif extension == "pptx":
        pages = extract_text_from_pptx(file_path)
        
    else:
        print(f"❌ Unsupported file type: {extension}")
        return []

    if not pages:
        print(f"⚠️ No text could be extracted from {file_name}")
        return []

    # --- STEP 2: Splitting Logic ---
    # We use a slightly larger chunk for technical OS concepts
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=600,
        chunk_overlap=120,
        separators=["\n\n", "\n", ". ", " ", ""]
    )

    all_chunks = []

    for page in pages:
        chunks = splitter.split_text(page["text"])

        for chunk in chunks:
            if len(chunk.strip()) < 40:
                continue

            doc = Document(
                page_content=chunk,
                metadata={
                    "source_file": file_name,
                    "doc_type": doc_type,
                    "subject": subject,
                    "page_number": page.get("page_number", 1),
                    "file_path": file_path
                }
            )
            all_chunks.append(doc)

    print(f"✅ {file_name} → {len(all_chunks)} chunks (Type: {doc_type})")
    return all_chunks