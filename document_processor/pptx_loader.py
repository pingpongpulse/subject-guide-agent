from pptx import Presentation
from pptx.util import Inches
import os

def load_pptx(filepath):
    """
    Extracts text from PPTX files including:
    - Text boxes and shapes
    - Tables inside slides
    - Speaker notes
    - Slide titles
    """
    prs = Presentation(filepath)
    slides = []

    for slide_num, slide in enumerate(prs.slides):
        slide_text = []

        for shape in slide.shapes:

            # Extract text from normal shapes and text boxes
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            # Extract text from tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        [cell.text.strip() for cell in row.cells
                         if cell.text.strip()]
                    )
                    if row_text:
                        slide_text.append(row_text)

        # Extract speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                slide_text.append(f"[Notes]: {notes.strip()}")

        if slide_text:
            slides.append({
                "text": "\n".join(slide_text),
                "page_number": slide_num + 1,
                "source_file": os.path.basename(filepath),
                "total_pages": len(prs.slides),
                "extraction_method": "python-pptx"
            })

    return slides


if __name__ == "__main__":
    result = load_pptx("data/sample_docs/test.pptx")
    for slide in result:
        print(f"--- Slide {slide['page_number']} ---")
        print(slide['text'][:300])
        print()